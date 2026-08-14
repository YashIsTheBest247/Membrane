"""Deck design system.

Register: an editorial minimal template. Pure monochrome — black, white and
three greys, no accent colour anywhere. Hierarchy comes from scale and space:
an 8pt letterspaced label above a 38pt caps headline, hairline rules, and a
great deal of nothing.

The rules this file enforces:
  · no colour, no gradients, no shadows, no fills except black and white
  · one family of rules, all 0.75pt
  · labels are letterspaced caps; headlines are caps; body is sentence case
  · content never exceeds two thirds of the slide — the rest is air
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
SHOTS = HERE.parent / "shots"
DIAGRAMS = HERE.parent / "diagrams"
OUT = HERE / "Membrane_NullDeity_NGH26_132.pptx"

W = Inches(13.333)
H = Inches(7.5)

# --------------------------------------------------------------------------
# palette — black, white, three greys. Nothing else.
# --------------------------------------------------------------------------

INK = RGBColor(0x11, 0x11, 0x11)      # headlines, black slides
GREY_DARK = RGBColor(0x44, 0x44, 0x44)  # body on white
GREY = RGBColor(0x77, 0x77, 0x77)     # secondary
GREY_LIGHT = RGBColor(0xA6, 0xA6, 0xA4)  # labels, page numbers
RULE = RGBColor(0xDD, 0xDD, 0xDA)     # hairlines
PAPER = RGBColor(0xF7, 0xF7, 0xF5)    # the ground
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Two accents, used sparingly and always meaningfully. ACCENT is the product's
# own brand indigo, taken from the logo so the deck and the software agree;
# DANGER marks the attack side of a contrast. Nothing else in the deck is
# coloured, so either one reads instantly.
ACCENT = RGBColor(0x5B, 0x52, 0xE0)
ACCENT_SOFT = RGBColor(0xE9, 0xE7, 0xFB)
DANGER = RGBColor(0xD8, 0x45, 0x3A)
DANGER_SOFT = RGBColor(0xFA, 0xEA, 0xE8)

# On black panels
INK_BG = RGBColor(0x14, 0x14, 0x14)
ON_INK = RGBColor(0xF4, 0xF4, 0xF2)
ON_INK_GREY = RGBColor(0x9A, 0x9A, 0x98)
ON_INK_RULE = RGBColor(0x33, 0x33, 0x33)

# Bahnschrift is the Windows geometric grotesque and matches the reference;
# PowerPoint falls back gracefully where it is absent.
DISPLAY = "Bahnschrift SemiBold"
DISPLAY_LIGHT = "Bahnschrift Light"
BODY = "Segoe UI"
MONO = "Consolas"

MARGIN = 1.0            # left and right margin, inches
TEAM = "Team NullDeity · NGH26_132"
TRACK = "Track 02 · Cybersecurity for the Future"


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    box.text_frame.word_wrap = True
    return box


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, font=BODY, spacing=None):
    """runs: (string, size, colour, bold[, font]) or a list of those.

    `spacing` is letter spacing in points — the label look depends on it, and
    python-pptx has no API for it, so it is set on the run's rPr directly.
    """
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    if isinstance(runs, tuple):
        runs = [runs]

    for index, item in enumerate(runs):
        content, size, colour, bold = item[0], item[1], item[2], item[3]
        font_name = item[4] if len(item) > 4 else font
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        run = para.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = font_name
        if spacing:
            run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


def image_fit(slide, path: Path, x, y, w, h, *, center=True):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    dx = x + (w - dw) // 2 if center else x
    dy = y + (h - dh) // 2 if center else y
    return slide.shapes.add_picture(str(path), dx, dy, dw, dh)


# --------------------------------------------------------------------------
# the type system
# --------------------------------------------------------------------------


def label(slide, x, y, content, *, dark=False, w=6.0, align=PP_ALIGN.LEFT):
    """The 8pt letterspaced grey caps that sit above every headline."""
    return text(slide, x, y, Inches(w), Inches(0.24),
                (content.upper(), 8, ON_INK_GREY if dark else ACCENT, False),
                spacing=2.6, align=align)


def headline(slide, x, y, content, *, dark=False, size=38, w=11.4,
             align=PP_ALIGN.LEFT):
    return text(slide, x, y, Inches(w), Inches(size / 52),
                (content.upper(), size, ON_INK if dark else INK, False),
                font=DISPLAY, spacing=0.4, align=align, line_spacing=1.02)


def body(slide, x, y, w, h, content, *, dark=False, size=11, colour=None,
         align=PP_ALIGN.LEFT):
    return text(slide, x, y, w, h,
                (content, size, colour or (ON_INK_GREY if dark else GREY_DARK), False),
                line_spacing=1.62, align=align)


def hairline(slide, x, y, w, *, dark=False, weight=0.012):
    return rect(slide, x, y, w, Inches(weight),
                fill=ON_INK_RULE if dark else RULE)


def numeral(slide, x, y, content, *, dark=False, size=13):
    """A small sequence numeral — 01, 02, 03 — set in the display face."""
    return text(slide, x, y, Inches(0.7), Inches(0.3),
                (content, size, ON_INK_GREY if dark else ACCENT, False),
                font=DISPLAY, spacing=0.8)


def outline_tag(slide, x, y, w, h, content, *, dark=False):
    """The small outlined caps box the reference uses for calls to action."""
    rect(slide, x, y, w, h, fill=None,
         line=ON_INK_RULE if dark else RGBColor(0xC8, 0xC8, 0xC5))
    text(slide, x, y + Inches(0.085), w, Inches(0.24),
         (content.upper(), 8, ON_INK if dark else INK, False),
         spacing=2.2, align=PP_ALIGN.CENTER)


def ring(slide, cx, cy, diameter, fraction, caption, value, *, dark=False):
    """A thin ring gauge: light track, dark arc, value set inside.

    Drawn as two stroked arcs rather than filled donuts so the line weight
    matches the hairlines everywhere else.
    """
    size = Inches(diameter)
    x = Inches(cx - diameter / 2)
    y = Inches(cy - diameter / 2)

    track = slide.shapes.add_shape(MSO_SHAPE.ARC, x, y, size, size)
    track.shadow.inherit = False
    track.fill.background()
    track.line.color.rgb = ON_INK_RULE if dark else RULE
    track.line.width = Pt(2.5)
    track.adjustments[0] = -90.0
    track.adjustments[1] = 269.9

    sweep = max(4.0, 360.0 * max(0.0, min(1.0, fraction)))
    arc = slide.shapes.add_shape(MSO_SHAPE.ARC, x, y, size, size)
    arc.shadow.inherit = False
    arc.fill.background()
    arc.line.color.rgb = ON_INK if dark else ACCENT
    arc.line.width = Pt(3.0)
    arc.adjustments[0] = -90.0
    arc.adjustments[1] = -90.0 + sweep

    text(slide, x, Inches(cy - 0.24), size, Inches(0.44),
         (value, 20, ON_INK if dark else INK, False),
         font=DISPLAY, align=PP_ALIGN.CENTER)
    text(slide, Inches(cx - diameter / 2 - 0.5), Inches(cy + diameter / 2 + 0.16),
         Inches(diameter + 1.0), Inches(0.26),
         (caption.upper(), 8, ON_INK_GREY if dark else GREY_LIGHT, False),
         spacing=2.0, align=PP_ALIGN.CENTER)


# --------------------------------------------------------------------------
# grounds and chrome
# --------------------------------------------------------------------------


def white_bg(slide):
    rect(slide, 0, 0, W, H, fill=PAPER)


def black_bg(slide):
    rect(slide, 0, 0, W, H, fill=INK_BG)


def slide_head(slide, kicker: str, title: str, *, dark=False, sub: str | None = None,
               sub_w=7.6):
    """Label, headline, hairline — the masthead every content slide shares."""
    x = Inches(MARGIN)
    label(slide, x, Inches(0.86), kicker, dark=dark)
    headline(slide, x, Inches(1.16), title, dark=dark)
    hairline(slide, x, Inches(1.92), W - Inches(MARGIN * 2), dark=dark)
    if sub:
        body(slide, x, Inches(2.12), Inches(sub_w), Inches(0.6), sub,
             dark=dark, size=11.5,
             colour=ON_INK_GREY if dark else GREY)


def footer(slide, number: int, *, dark=False):
    colour = ON_INK_RULE if dark else GREY_LIGHT
    text(slide, Inches(MARGIN), Inches(6.96), Inches(6.0), Inches(0.26),
         ("MEMBRANE", 8, colour, False), spacing=2.4)
    text(slide, W - Inches(MARGIN + 1.0), Inches(6.96), Inches(1.0), Inches(0.26),
         (f"{number:02d}", 8, colour, False), spacing=2.4, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------
# composite blocks
# --------------------------------------------------------------------------


def listing(slide, x, y, w, items, *, dark=False, size=11, gap=0.74, numbered=False,
            start=1):
    """A list set as hairline-separated rows, optionally numbered."""
    top = y
    for index, item in enumerate(items):
        if numbered:
            numeral(slide, x, top, f"{start + index:02d}", dark=dark, size=11)
            tx, tw = x + Inches(0.62), w - Inches(0.62)
        else:
            tx, tw = x, w

        if isinstance(item, tuple):
            lead, rest = item
            box = slide.shapes.add_textbox(tx, top, tw, Inches(0.4))
            frame = box.text_frame
            frame.word_wrap = True
            frame.margin_left = frame.margin_right = 0
            frame.margin_top = frame.margin_bottom = 0
            para = frame.paragraphs[0]
            para.line_spacing = 1.6
            r1 = para.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = False
            r1.font.color.rgb = ON_INK if dark else INK
            r1.font.name = BODY
            if rest:
                r2 = para.add_run()
                r2.text = rest
                r2.font.size = Pt(size)
                r2.font.color.rgb = ON_INK_GREY if dark else GREY_DARK
                r2.font.name = BODY
        else:
            body(slide, tx, top, tw, Inches(0.4), item, dark=dark, size=size)
        top += Inches(gap)
    return top


def column(slide, x, y, w, kicker, title, content, *, dark=False, rule_above=True,
           title_size=15):
    """A titled column: hairline, label, title, body."""
    if rule_above:
        hairline(slide, x, y, w, dark=dark)
    label(slide, x, y + Inches(0.22), kicker, dark=dark, w=w / 914400)
    text(slide, x, y + Inches(0.52), w, Inches(0.36),
         (title, title_size, ON_INK if dark else INK, False), font=DISPLAY,
         spacing=0.2)
    body(slide, x, y + Inches(1.0), w, Inches(1.5), content, dark=dark)


def table(slide, x, y, w, headers, rows, *, col_w, dark=False, row_h=0.46,
          size=10.5):
    """Horizontal hairlines only. No fills, no vertical rules."""
    cx = x
    for index, header in enumerate(headers):
        text(slide, cx, y, Inches(col_w[index] - 0.25), Inches(0.24),
             (header.upper(), 8, ON_INK_GREY if dark else GREY_LIGHT, False),
             spacing=2.2)
        cx += Inches(col_w[index])

    ty = y + Inches(0.34)
    hairline(slide, x, ty, w, dark=dark)
    ty += Inches(0.02)

    for row in rows:
        cx = x
        for index, cell in enumerate(row):
            first = index == 0
            text(slide, cx, ty + Inches(0.16), Inches(col_w[index] - 0.3),
                 Inches(row_h - 0.12),
                 (cell, size, (ON_INK if dark else INK) if first
                  else (ON_INK_GREY if dark else GREY_DARK), False),
                 line_spacing=1.5)
            cx += Inches(col_w[index])
        ty += Inches(row_h)
        hairline(slide, x, ty, w, dark=dark)
        ty += Inches(0.02)
    return ty


# --------------------------------------------------------------------------
# specimen blocks — the deck's "photography"
# --------------------------------------------------------------------------


def specimen(slide, x, y, w, h, lines, *, dark=True, size=9.5, caption=None):
    """A monospace specimen of real content.

    This deck has no stock photography, and it should not pretend to. What it
    has instead is the actual material — a poisoned page, a verdict, a
    provenance chain — set as a block. It carries the visual weight an image
    would, and unlike an image it is evidence.

    `lines` is a list of (text, emphasised) pairs; emphasised lines are set in
    full contrast, the rest recede.
    """
    rect(slide, x, y, w, h, fill=INK_BG if dark else WHITE,
         line=None if dark else RULE)

    box = slide.shapes.add_textbox(x + Inches(0.32), y + Inches(0.28),
                                   w - Inches(0.64), h - Inches(0.5))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    for index, (content, strong) in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.line_spacing = 1.5
        run = para.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.name = MONO
        if dark:
            run.font.color.rgb = (RGBColor(0xFF, 0x8A, 0x7A) if strong
                                  else RGBColor(0x8A, 0x8A, 0x88))
        else:
            run.font.color.rgb = DANGER if strong else GREY

    if caption:
        text(slide, x, y + h + Inches(0.14), w, Inches(0.24),
             (caption.upper(), 8, GREY_LIGHT, False), spacing=2.2)
    return box


def display_number(slide, x, y, w, content, *, dark=False, size=104,
                   align=PP_ALIGN.LEFT, colour=None):
    """An oversized numeral, used where a figure is the whole point."""
    return text(slide, x, y, w, Inches(size / 58),
                (content, size, colour or (ON_INK if dark else INK), False),
                font=DISPLAY, align=align, line_spacing=0.94, spacing=-1.0)
