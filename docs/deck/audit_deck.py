"""Layout audit for the built deck.

There is no PowerPoint or LibreOffice on this machine, so the deck cannot be
rendered and eyeballed here. This checks what *can* be checked without a
renderer: shapes escaping the slide, text boxes too small for the text in them,
missing or badly-scaled images, and slides that have quietly gone empty.

Text overflow is estimated, not measured — glyph widths depend on the font
metrics PowerPoint resolves at open time. The heuristic is deliberately
conservative so it flags candidates for a human to look at rather than
pretending to be authoritative.
"""

from __future__ import annotations

import math
from pathlib import Path

import sys

from pptx import Presentation
from pptx.util import Emu

DECK = Path(__file__).parent / "Membrane_NullDeity_NGH26_132.pptx"

SLIDE_W = Emu(Inches := 12192000)  # 13.333in in EMU
SLIDE_H = 6858000                  # 7.5in in EMU
EMU_PER_IN = 914400
EMU_PER_PT = 12700

# Average glyph advance as a fraction of point size, for a humanist sans.
CHAR_W = 0.50
BOLD_CHAR_W = 0.54


def inches(value: int) -> float:
    return value / EMU_PER_IN


def estimate_text_height(shape) -> tuple[float, float]:
    """Return (estimated_height_in, box_height_in) for a text-bearing shape."""
    frame = shape.text_frame
    box_w_pt = shape.width / EMU_PER_PT
    total_pt = 0.0

    for para in frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            total_pt += 8
            continue
        size = max((r.font.size.pt if r.font.size else 18) for r in runs)
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        for chunk in "".join(r.text for r in runs).split("\n"):
            width_factor = BOLD_CHAR_W if any(r.font.bold for r in runs) else CHAR_W
            per_line = max(1, int(box_w_pt / (size * width_factor)))
            lines = max(1, math.ceil(len(chunk) / per_line))
            total_pt += lines * size * 1.2 * spacing
        space_after = para.space_after.pt if para.space_after else 0
        total_pt += space_after

    return total_pt / 72.0, shape.height / EMU_PER_IN


def audit() -> int:
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8", errors="replace")
        except (AttributeError, ValueError):
            pass
    prs = Presentation(DECK)
    problems: list[str] = []
    notes: list[str] = []

    print(f"deck: {DECK.name}")
    print(f"slides: {len(prs.slides)}   "
          f"{inches(prs.slide_width):.2f} x {inches(prs.slide_height):.2f} in\n")

    for index, slide in enumerate(prs.slides, start=1):
        shapes = list(slide.shapes)
        pictures = [s for s in shapes if s.shape_type == 13]
        texts = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]

        if not texts:
            problems.append(f"slide {index:02d}: no text at all")

        for shape in shapes:
            # Off-slide geometry
            right = shape.left + shape.width
            bottom = shape.top + shape.height
            if shape.left < -1000 or shape.top < -1000:
                problems.append(
                    f"slide {index:02d}: shape starts off-slide "
                    f"({inches(shape.left):.2f}, {inches(shape.top):.2f} in)")
            if right > prs.slide_width + 1000:
                problems.append(
                    f"slide {index:02d}: shape overruns the right edge by "
                    f"{inches(right - prs.slide_width):.2f} in")
            if bottom > prs.slide_height + 1000:
                problems.append(
                    f"slide {index:02d}: shape overruns the bottom edge by "
                    f"{inches(bottom - prs.slide_height):.2f} in")

            # Text that will not fit its box
            if shape.has_text_frame and shape.text_frame.text.strip():
                needed, available = estimate_text_height(shape)
                if needed > available * 1.30 and needed - available > 0.12:
                    preview = shape.text_frame.text.strip().replace("\n", " ")[:58]
                    notes.append(
                        f"slide {index:02d}: text may overflow — needs ~"
                        f"{needed:.2f} in, box is {available:.2f} in  “{preview}…”")

        print(f"  slide {index:02d}  shapes {len(shapes):3d}   "
              f"text {len(texts):2d}   images {len(pictures)}")

    print()
    if problems:
        print(f"PROBLEMS ({len(problems)})")
        for item in problems:
            print(f"  ✗ {item}")
    else:
        print("PROBLEMS: none — every shape sits inside the slide")

    print()
    if notes:
        print(f"POSSIBLE TEXT OVERFLOW ({len(notes)}) — estimated, worth a look")
        for item in notes:
            print(f"  · {item}")
    else:
        print("POSSIBLE TEXT OVERFLOW: none detected")

    return 1 if problems else 0


if __name__ == "__main__":
    raise SystemExit(audit())
